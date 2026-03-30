from __future__ import annotations

import base64
import csv
import io
import zipfile
from collections import deque
from dataclasses import dataclass, field
from pathlib import PurePosixPath
from typing import Any
from urllib.parse import quote
from uuid import UUID, uuid4

import httpx
import pytesseract
from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image
from pypdf import PdfReader
from pptx import Presentation
from sqlalchemy import func, select
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker

from app.core.config import get_settings
from app.core.security import (
    create_code_challenge,
    decrypt_secret,
    encrypt_secret,
    future_utc,
    generate_code_verifier,
    generate_state_token,
)
from app.core.utils import normalize_whitespace, slugify, utcnow
from app.db.models import (
    ConnectorConnection,
    ConnectorOAuthPurpose,
    ConnectorOAuthState,
    ConnectorOwnerScope,
    ConnectorProvider,
    ConnectorResource,
    ConnectorResourceKind,
    ConnectorResourceStatus,
    ConnectorSourceItem,
    ConnectorSourceItemStatus,
    ConnectorStatus,
    ConnectorSyncJob,
    ConnectorSyncJobKind,
    ConnectorSyncMode,
    Document,
    DocumentVisibilityScope,
    JobStatus,
)
from app.schemas.connectors import (
    ConnectorBrowseItem,
    ConnectorBrowseResponse,
    ConnectorConnectionSummary,
    ConnectorListResponse,
    ConnectorOAuthCallbackResponse,
    ConnectorProviderReadiness,
    ConnectorReadinessResponse,
    ConnectorResourceCreateRequest,
    ConnectorResourceSummary,
    ConnectorResourceUpdateRequest,
    ConnectorSourceItemSummary,
    ConnectorUpdateRequest,
)
from app.schemas.documents import IngestDocumentRequest
from app.schemas.jobs import JobSummary
from app.services.auth import AuthenticatedUser
from app.services.ingest import ingest_document
from app.services.parser import DocumentParser

GOOGLE_AUTH_URL = "https://accounts.google.com/o/oauth2/v2/auth"
GOOGLE_TOKEN_URL = "https://oauth2.googleapis.com/token"
GOOGLE_USERINFO_URL = "https://openidconnect.googleapis.com/v1/userinfo"
GOOGLE_DRIVE_V3 = "https://www.googleapis.com/drive/v3"
GOOGLE_CONNECTOR_SCOPES = ["openid", "email", "profile", "https://www.googleapis.com/auth/drive.readonly"]
GOOGLE_FOLDER_MIME = "application/vnd.google-apps.folder"
GOOGLE_DOC_MIME = "application/vnd.google-apps.document"
GOOGLE_SHEET_MIME = "application/vnd.google-apps.spreadsheet"
GOOGLE_SLIDE_MIME = "application/vnd.google-apps.presentation"
GOOGLE_DRAWING_MIME = "application/vnd.google-apps.drawing"

NOTION_AUTH_URL = "https://api.notion.com/v1/oauth/authorize"
NOTION_TOKEN_URL = "https://api.notion.com/v1/oauth/token"
NOTION_API_URL = "https://api.notion.com/v1"
NOTION_VERSION = "2022-06-28"

GITHUB_AUTH_URL = "https://github.com/login/oauth/authorize"
GITHUB_TOKEN_URL = "https://github.com/login/oauth/access_token"
GITHUB_API_URL = "https://api.github.com"
GITHUB_CONNECTOR_SCOPES = ["read:user", "user:email", "repo"]
GITHUB_DOC_EXTENSIONS = {".md", ".mdx", ".rst", ".adoc", ".txt"}
NOTION_EXPORT_EXTENSIONS = {".md", ".markdown", ".html", ".htm", ".txt", ".csv"}

SYNC_INTERVALS = {15, 60, 360, 1440}
PROVIDER_PATH_TO_VALUE = {
    "google-drive": ConnectorProvider.google_drive.value,
    "google_drive": ConnectorProvider.google_drive.value,
    "github": ConnectorProvider.github.value,
    "notion": ConnectorProvider.notion.value,
}
PROVIDER_VALUE_TO_PATH = {
    ConnectorProvider.google_drive.value: "google-drive",
    ConnectorProvider.github.value: "github",
    ConnectorProvider.notion.value: "notion",
}
RESOURCE_KINDS_BY_PROVIDER = {
    ConnectorProvider.google_drive.value: {
        ConnectorResourceKind.folder.value,
        ConnectorResourceKind.shared_drive.value,
    },
    ConnectorProvider.github.value: {
        ConnectorResourceKind.repository_docs.value,
        ConnectorResourceKind.repository_evidence.value,
    },
    ConnectorProvider.notion.value: {
        ConnectorResourceKind.page.value,
        ConnectorResourceKind.database.value,
        ConnectorResourceKind.export_upload.value,
    },
}
RECOMMENDED_TEMPLATES_BY_PROVIDER = {
    ConnectorProvider.google_drive.value: ["shared_drive", "folder"],
    ConnectorProvider.github.value: ["repository_docs", "repository_evidence"],
    ConnectorProvider.notion.value: ["page", "database", "export_upload"],
}


class ConnectorError(RuntimeError):
    pass


class ConnectorForbiddenError(ConnectorError):
    pass


class ConnectorNotFoundError(ConnectorError):
    pass


@dataclass(slots=True)
class ExtractedContent:
    content_type: str
    content: str
    doc_type: str


@dataclass(slots=True)
class PreparedSyncItem:
    external_item_id: str
    title: str
    source_url: str | None
    source_revision_id: str | None
    mime_type: str | None
    content_type: str | None
    content: str | None
    doc_type: str
    unsupported_reason: str | None = None
    provider_metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(slots=True)
class ProviderOAuthResult:
    account_subject: str
    display_name: str
    account_email: str | None
    encrypted_access_token: str
    encrypted_refresh_token: str | None
    token_expires_at: Any
    granted_scopes: list[str]


def _app_callback_path(path: str) -> str:
    settings = get_settings()
    return f"{settings.app_public_url.rstrip('/')}{path}"


def _safe_return_path(value: str | None) -> str:
    if not value or not value.startswith("/"):
        return "/connectors"
    return value


def _normalize_provider(value: str) -> str:
    provider = PROVIDER_PATH_TO_VALUE.get(value)
    if provider is None:
        raise ConnectorError("Unsupported connector provider.")
    return provider


def _provider_path(provider: str) -> str:
    return PROVIDER_VALUE_TO_PATH.get(provider, provider)


def _validate_owner_scope(scope: str) -> str:
    normalized = scope.strip().lower()
    if normalized == "shared":
        normalized = ConnectorOwnerScope.workspace.value
    elif normalized == "user":
        normalized = ConnectorOwnerScope.personal.value
    if normalized not in {ConnectorOwnerScope.workspace.value, ConnectorOwnerScope.personal.value}:
        raise ConnectorError("Unsupported owner scope.")
    return normalized


def _validate_resource_kind(provider: str, value: str | None) -> str:
    if not value:
        return (
            ConnectorResourceKind.folder.value
            if provider == ConnectorProvider.google_drive.value
            else ConnectorResourceKind.repository_docs.value
            if provider == ConnectorProvider.github.value
            else ConnectorResourceKind.page.value
        )
    if value not in RESOURCE_KINDS_BY_PROVIDER[provider]:
        raise ConnectorError("Unsupported connector resource kind.")
    return value


def _validate_browse_kind(provider: str, value: str | None) -> str:
    if provider == ConnectorProvider.notion.value and value == ConnectorResourceKind.export_upload.value:
        raise ConnectorError("Uploaded exports do not support browse.")
    return _validate_resource_kind(provider, value)


def _default_visibility_scope(provider: str, resource_kind: str) -> str:
    if provider == ConnectorProvider.github.value and resource_kind == ConnectorResourceKind.repository_evidence.value:
        return DocumentVisibilityScope.evidence_only.value
    if provider == ConnectorProvider.notion.value and resource_kind == ConnectorResourceKind.export_upload.value:
        return DocumentVisibilityScope.evidence_only.value
    return DocumentVisibilityScope.member_visible.value


def _normalize_visibility_scope(provider: str, resource_kind: str, value: str | None) -> str:
    normalized = (value or _default_visibility_scope(provider, resource_kind)).strip().lower()
    if normalized not in {
        DocumentVisibilityScope.member_visible.value,
        DocumentVisibilityScope.evidence_only.value,
    }:
        raise ConnectorError("Unsupported visibility scope.")
    return normalized


def _default_selection_mode(provider: str, resource_kind: str) -> str:
    if provider == ConnectorProvider.notion.value and resource_kind == ConnectorResourceKind.export_upload.value:
        return "export_upload"
    if provider == ConnectorProvider.github.value:
        return "search"
    if provider == ConnectorProvider.notion.value:
        return "search"
    return "browse"


def _normalize_selection_mode(provider: str, resource_kind: str, value: str | None) -> str:
    normalized = (value or _default_selection_mode(provider, resource_kind)).strip().lower()
    if not normalized:
        raise ConnectorError("Selection mode is required.")
    return normalized


def _normalize_sync_schedule(sync_mode: str, interval: int | None) -> tuple[str, int | None]:
    if sync_mode not in {ConnectorSyncMode.manual.value, ConnectorSyncMode.auto.value}:
        raise ConnectorError("Unsupported sync mode.")
    if sync_mode == ConnectorSyncMode.auto.value:
        normalized_interval = interval if interval in SYNC_INTERVALS else 60
        return sync_mode, normalized_interval
    return ConnectorSyncMode.manual.value, None


def _default_sync_schedule_for_scope(owner_scope: str) -> tuple[str, int | None]:
    if owner_scope == ConnectorOwnerScope.workspace.value:
        return ConnectorSyncMode.auto.value, 60
    return ConnectorSyncMode.manual.value, None


def _resource_supports_connector_sync(resource: ConnectorResource) -> bool:
    return resource.selection_mode != "export_upload"


def _non_syncable_resource_message(resource: ConnectorResource) -> str:
    if resource.selection_mode == "export_upload":
        return "업로드형 내보내기는 새 파일 업로드로 갱신하세요."
    return "이 소스는 커넥터 동기화를 지원하지 않습니다."


def _resource_sync_defaults(connection: ConnectorConnection, payload: ConnectorResourceCreateRequest) -> tuple[bool, str, int | None]:
    provider = connection.provider
    kind = _validate_resource_kind(provider, payload.resource_kind)
    if provider == ConnectorProvider.notion.value and kind == ConnectorResourceKind.export_upload.value:
        return False, ConnectorSyncMode.manual.value, None
    default_sync_mode, default_interval = _default_sync_schedule_for_scope(connection.owner_scope)
    sync_mode, interval = _normalize_sync_schedule(
        payload.sync_mode or default_sync_mode,
        payload.sync_interval_minutes if payload.sync_interval_minutes is not None else default_interval,
    )
    if provider == ConnectorProvider.google_drive.value:
        sync_children = True if payload.sync_children is None else payload.sync_children
        return sync_children, sync_mode, interval
    if kind == ConnectorResourceKind.database.value:
        return True, sync_mode, interval
    return False, sync_mode, interval


def _resource_sync_children_for_update(resource: ConnectorResource, requested: bool | None) -> bool:
    if resource.provider == ConnectorProvider.google_drive.value:
        return resource.sync_children if requested is None else requested
    if resource.resource_kind == ConnectorResourceKind.database.value:
        return True
    return False


def _stable_document_slug(title: str, external_id: str) -> str:
    return f"{slugify(title)}-{external_id[:10].lower()}"


def _google_revision_token(file_meta: dict[str, Any]) -> str:
    version = str(file_meta.get("version") or "")
    modified = str(file_meta.get("modifiedTime") or "")
    checksum = str(file_meta.get("md5Checksum") or "")
    return f"{version}:{modified}:{checksum}"


def _notion_revision_token(page: dict[str, Any]) -> str:
    return str(page.get("last_edited_time") or "")


def _google_drive_configured() -> bool:
    settings = get_settings()
    return bool(settings.google_oauth_client_id and settings.google_oauth_client_secret)


def _github_configured() -> bool:
    settings = get_settings()
    return bool(settings.github_oauth_client_id and settings.github_oauth_client_secret)


def _notion_configured() -> bool:
    settings = get_settings()
    return bool(settings.notion_oauth_client_id and settings.notion_oauth_client_secret)


def _provider_configured(provider: str) -> bool:
    if provider == ConnectorProvider.google_drive.value:
        return _google_drive_configured()
    if provider == ConnectorProvider.github.value:
        return _github_configured()
    if provider == ConnectorProvider.notion.value:
        return _notion_configured()
    return False


def _ensure_provider_configured(provider: str) -> None:
    if not _provider_configured(provider):
        raise ConnectorError(f"{provider} OAuth is not configured.")


def _connection_summary(connection: ConnectorConnection, resources: list[ConnectorResource]) -> ConnectorConnectionSummary:
    return ConnectorConnectionSummary(
        id=connection.id,
        provider=connection.provider,
        owner_scope=connection.owner_scope,
        owner_user_id=connection.owner_user_id,
        display_name=connection.display_name,
        account_email=connection.account_email,
        account_subject=connection.account_subject,
        status=connection.status,
        granted_scopes=list(connection.granted_scopes or []),
        last_validated_at=connection.last_validated_at,
        created_at=connection.created_at,
        updated_at=connection.updated_at,
        resources=[_resource_summary(resource) for resource in resources],
    )


def _resource_summary(resource: ConnectorResource) -> ConnectorResourceSummary:
    return ConnectorResourceSummary(
        id=resource.id,
        connection_id=resource.connection_id,
        provider=resource.provider,
        resource_kind=resource.resource_kind,
        external_id=resource.external_id,
        name=resource.name,
        resource_url=resource.resource_url,
        parent_external_id=resource.parent_external_id,
        visibility_scope=resource.visibility_scope,
        selection_mode=resource.selection_mode,
        sync_children=resource.sync_children,
        sync_mode=resource.sync_mode,
        sync_interval_minutes=resource.sync_interval_minutes,
        status=resource.status,
        last_sync_started_at=resource.last_sync_started_at,
        last_sync_completed_at=resource.last_sync_completed_at,
        next_auto_sync_at=resource.next_auto_sync_at,
        last_sync_summary=dict(resource.last_sync_summary or {}),
        provider_metadata=dict(resource.provider_metadata or {}),
    )


def _source_item_summary(item: ConnectorSourceItem) -> ConnectorSourceItemSummary:
    return ConnectorSourceItemSummary(
        id=item.id,
        resource_id=item.resource_id,
        external_item_id=item.external_item_id,
        mime_type=item.mime_type,
        name=item.name,
        source_url=item.source_url,
        source_revision_id=item.source_revision_id,
        internal_document_id=item.internal_document_id,
        item_status=item.item_status,
        unsupported_reason=item.unsupported_reason,
        error_message=item.error_message,
        last_seen_at=item.last_seen_at,
        last_synced_at=item.last_synced_at,
        provider_metadata=dict(item.provider_metadata or {}),
    )


def _provider_readiness_summary(
    provider: str,
    *,
    connection: ConnectorConnection | None,
    auth_user: AuthenticatedUser | None,
    healthy_source_count: int = 0,
    needs_attention_count: int = 0,
) -> ConnectorProviderReadiness:
    oauth_configured = _provider_configured(provider)
    if not oauth_configured:
        setup_state = "not_configured"
    elif connection is None:
        setup_state = "setup_needed"
    elif needs_attention_count > 0 or connection.status != ConnectorStatus.active.value:
        setup_state = "attention_required"
    else:
        setup_state = "ready"
    return ConnectorProviderReadiness(
        provider=provider,
        oauth_configured=oauth_configured,
        workspace_connection_exists=connection is not None,
        workspace_connection_status=connection.status if connection is not None else None,
        viewer_can_manage_workspace_connection=bool(auth_user and auth_user.can_manage_workspace_connectors),
        setup_state=setup_state,
        healthy_source_count=healthy_source_count,
        needs_attention_count=needs_attention_count,
        recommended_templates=list(RECOMMENDED_TEMPLATES_BY_PROVIDER.get(provider, [])),
    )


async def _exchange_google_code(*, code: str, code_verifier: str, redirect_uri: str) -> dict[str, Any]:
    _ensure_provider_configured(ConnectorProvider.google_drive.value)
    settings = get_settings()
    payload = {
        "client_id": settings.google_oauth_client_id,
        "client_secret": settings.google_oauth_client_secret,
        "code": code,
        "code_verifier": code_verifier,
        "grant_type": "authorization_code",
        "redirect_uri": redirect_uri,
    }
    async with httpx.AsyncClient(timeout=30) as client:
        response = await client.post(GOOGLE_TOKEN_URL, data=payload)
    if response.status_code >= 400:
        raise ConnectorError(f"Google token exchange failed: {response.text}")
    return response.json()


async def _refresh_google_token(refresh_token: str) -> dict[str, Any]:
    _ensure_provider_configured(ConnectorProvider.google_drive.value)
    settings = get_settings()
    payload = {
        "client_id": settings.google_oauth_client_id,
        "client_secret": settings.google_oauth_client_secret,
        "refresh_token": refresh_token,
        "grant_type": "refresh_token",
    }
    async with httpx.AsyncClient(timeout=30) as client:
        response = await client.post(GOOGLE_TOKEN_URL, data=payload)
    if response.status_code >= 400:
        raise ConnectorError(f"Google token refresh failed: {response.text}")
    return response.json()


async def _google_userinfo(access_token: str) -> dict[str, Any]:
    async with httpx.AsyncClient(timeout=30) as client:
        response = await client.get(
            GOOGLE_USERINFO_URL,
            headers={"Authorization": f"Bearer {access_token}"},
        )
    if response.status_code >= 400:
        raise ConnectorError(f"Google userinfo lookup failed: {response.text}")
    return response.json()


async def _exchange_github_code(*, code: str, redirect_uri: str) -> dict[str, Any]:
    _ensure_provider_configured(ConnectorProvider.github.value)
    settings = get_settings()
    async with httpx.AsyncClient(timeout=30) as client:
        response = await client.post(
            GITHUB_TOKEN_URL,
            headers={
                "Accept": "application/json",
            },
            data={
                "client_id": settings.github_oauth_client_id,
                "client_secret": settings.github_oauth_client_secret,
                "code": code,
                "redirect_uri": redirect_uri,
            },
        )
    if response.status_code >= 400:
        raise ConnectorError(f"GitHub token exchange failed: {response.text}")
    payload = response.json()
    if payload.get("error"):
        raise ConnectorError(f"GitHub token exchange failed: {payload.get('error_description') or payload['error']}")
    return payload


async def _github_user(access_token: str) -> dict[str, Any]:
    async with httpx.AsyncClient(timeout=30) as client:
        response = await client.get(
            f"{GITHUB_API_URL}/user",
            headers={
                "Authorization": f"Bearer {access_token}",
                "Accept": "application/vnd.github+json",
                "X-GitHub-Api-Version": "2022-11-28",
            },
        )
    if response.status_code >= 400:
        raise ConnectorError(f"GitHub user lookup failed: {response.text}")
    return response.json()


async def _github_primary_email(access_token: str) -> str | None:
    async with httpx.AsyncClient(timeout=30) as client:
        response = await client.get(
            f"{GITHUB_API_URL}/user/emails",
            headers={
                "Authorization": f"Bearer {access_token}",
                "Accept": "application/vnd.github+json",
                "X-GitHub-Api-Version": "2022-11-28",
            },
        )
    if response.status_code >= 400:
        return None
    emails = response.json()
    if not isinstance(emails, list):
        return None
    for email in emails:
        if email.get("primary") and email.get("verified") and email.get("email"):
            return str(email["email"])
    for email in emails:
        if email.get("verified") and email.get("email"):
            return str(email["email"])
    return None


async def _exchange_notion_code(*, code: str, redirect_uri: str) -> dict[str, Any]:
    _ensure_provider_configured(ConnectorProvider.notion.value)
    settings = get_settings()
    basic = base64.b64encode(f"{settings.notion_oauth_client_id}:{settings.notion_oauth_client_secret}".encode("utf-8")).decode("utf-8")
    async with httpx.AsyncClient(timeout=30) as client:
        response = await client.post(
            NOTION_TOKEN_URL,
            headers={
                "Authorization": f"Basic {basic}",
                "Content-Type": "application/json",
            },
            json={
                "grant_type": "authorization_code",
                "code": code,
                "redirect_uri": redirect_uri,
            },
        )
    if response.status_code >= 400:
        raise ConnectorError(f"Notion token exchange failed: {response.text}")
    return response.json()


async def _active_access_token(session: AsyncSession, connection: ConnectorConnection) -> str:
    access_token = decrypt_secret(connection.encrypted_access_token)
    if not access_token:
        connection.status = ConnectorStatus.needs_reauth.value
        await session.commit()
        raise ConnectorError("Connector access token is unavailable.")
    if connection.provider in {ConnectorProvider.github.value, ConnectorProvider.notion.value}:
        return access_token
    refresh_token = decrypt_secret(connection.encrypted_refresh_token)
    expires_at = connection.token_expires_at
    if expires_at and expires_at > future_utc(seconds=60):
        return access_token
    if not refresh_token:
        connection.status = ConnectorStatus.needs_reauth.value
        await session.commit()
        raise ConnectorError("Connector token refresh is unavailable.")
    token_data = await _refresh_google_token(refresh_token)
    new_access = str(token_data["access_token"])
    connection.encrypted_access_token = encrypt_secret(new_access) or ""
    connection.token_expires_at = future_utc(seconds=int(token_data.get("expires_in", 3600)))
    if token_data.get("refresh_token"):
        connection.encrypted_refresh_token = encrypt_secret(str(token_data["refresh_token"]))
    connection.status = ConnectorStatus.active.value
    connection.last_validated_at = utcnow()
    await session.commit()
    return new_access


async def _google_json(
    session: AsyncSession,
    connection: ConnectorConnection,
    path: str,
    *,
    params: dict[str, Any] | None = None,
) -> dict[str, Any]:
    token = await _active_access_token(session, connection)
    async with httpx.AsyncClient(timeout=60) as client:
        response = await client.get(
            f"{GOOGLE_DRIVE_V3}{path}",
            params=params,
            headers={"Authorization": f"Bearer {token}"},
        )
    if response.status_code >= 400:
        raise ConnectorError(f"Google Drive request failed: {response.text}")
    return response.json()


async def _google_bytes(
    session: AsyncSession,
    connection: ConnectorConnection,
    path: str,
    *,
    params: dict[str, Any] | None = None,
) -> bytes:
    token = await _active_access_token(session, connection)
    async with httpx.AsyncClient(timeout=120) as client:
        response = await client.get(
            f"{GOOGLE_DRIVE_V3}{path}",
            params=params,
            headers={"Authorization": f"Bearer {token}"},
        )
    if response.status_code >= 400:
        raise ConnectorError(f"Google Drive download failed: {response.text}")
    return response.content


async def _github_request(
    session: AsyncSession,
    connection: ConnectorConnection,
    method: str,
    path: str,
    *,
    params: dict[str, Any] | None = None,
    expected_status: int = 200,
) -> tuple[Any, httpx.Headers]:
    token = await _active_access_token(session, connection)
    async with httpx.AsyncClient(timeout=60) as client:
        response = await client.request(
            method,
            f"{GITHUB_API_URL}{path}",
            params=params,
            headers={
                "Authorization": f"Bearer {token}",
                "Accept": "application/vnd.github+json",
                "X-GitHub-Api-Version": "2022-11-28",
            },
        )
    if response.status_code != expected_status:
        if response.status_code in {401, 403}:
            connection.status = ConnectorStatus.needs_reauth.value
            await session.commit()
        raise ConnectorError(f"GitHub request failed: {response.text}")
    return response.json(), response.headers


async def _github_raw_bytes(session: AsyncSession, connection: ConnectorConnection, download_url: str) -> bytes:
    token = await _active_access_token(session, connection)
    async with httpx.AsyncClient(timeout=60, follow_redirects=True) as client:
        response = await client.get(
            download_url,
            headers={
                "Authorization": f"Bearer {token}",
                "Accept": "application/vnd.github.raw",
            },
        )
    if response.status_code >= 400:
        raise ConnectorError(f"GitHub raw file download failed: {response.text}")
    return response.content


async def _notion_request(
    session: AsyncSession,
    connection: ConnectorConnection,
    method: str,
    path: str,
    *,
    params: dict[str, Any] | None = None,
    payload: dict[str, Any] | None = None,
) -> dict[str, Any]:
    token = await _active_access_token(session, connection)
    headers = {
        "Authorization": f"Bearer {token}",
        "Notion-Version": NOTION_VERSION,
    }
    if payload is not None:
        headers["Content-Type"] = "application/json"
    async with httpx.AsyncClient(timeout=60) as client:
        response = await client.request(
            method,
            f"{NOTION_API_URL}{path}",
            params=params,
            json=payload,
            headers=headers,
        )
    if response.status_code >= 400:
        if response.status_code in {401, 403}:
            connection.status = ConnectorStatus.needs_reauth.value
            await session.commit()
        raise ConnectorError(f"Notion request failed: {response.text}")
    return response.json()


def _notion_rich_text(rich_text: list[dict[str, Any]] | None) -> str:
    parts: list[str] = []
    for item in rich_text or []:
        plain = str(item.get("plain_text") or "")
        href = item.get("href")
        if href and plain:
            parts.append(f"[{plain}]({href})")
        else:
            parts.append(plain)
    return "".join(parts).strip()


def _notion_title_from_page(page: dict[str, Any]) -> str:
    for value in (page.get("properties") or {}).values():
        if isinstance(value, dict) and value.get("type") == "title":
            title = _notion_rich_text(value.get("title"))
            if title:
                return title
    return "Untitled"


def _notion_title_from_database(database: dict[str, Any]) -> str:
    return _notion_rich_text(database.get("title")) or "Untitled database"


async def _notion_block_children(session: AsyncSession, connection: ConnectorConnection, block_id: str) -> list[dict[str, Any]]:
    results: list[dict[str, Any]] = []
    cursor: str | None = None
    while True:
        params: dict[str, Any] = {"page_size": 100}
        if cursor:
            params["start_cursor"] = cursor
        payload = await _notion_request(session, connection, "GET", f"/blocks/{block_id}/children", params=params)
        results.extend(list(payload.get("results", [])))
        if not payload.get("has_more"):
            break
        cursor = payload.get("next_cursor")
        if not cursor:
            break
    return results


async def _notion_block_to_markdown(
    session: AsyncSession,
    connection: ConnectorConnection,
    block: dict[str, Any],
    *,
    depth: int = 0,
) -> str:
    block_type = str(block.get("type") or "")
    payload = block.get(block_type) or {}
    text = _notion_rich_text(payload.get("rich_text")) if isinstance(payload, dict) else ""
    prefix = ""
    if block_type == "heading_1":
        prefix = "# "
    elif block_type == "heading_2":
        prefix = "## "
    elif block_type == "heading_3":
        prefix = "### "
    elif block_type == "bulleted_list_item":
        prefix = "  " * depth + "- "
    elif block_type == "numbered_list_item":
        prefix = "  " * depth + "1. "
    elif block_type == "to_do":
        checked = payload.get("checked") is True
        prefix = "  " * depth + ("- [x] " if checked else "- [ ] ")
    elif block_type in {"quote", "callout"}:
        prefix = "> "
    elif block_type == "code":
        language = str(payload.get("language") or "")
        return f"```{language}\n{text}\n```".strip()
    elif block_type == "divider":
        return "---"
    elif block_type == "child_page":
        return f"## {payload.get('title') or text or 'Child page'}"
    elif block_type == "bookmark":
        url = payload.get("url")
        return f"[북마크]({url})" if url else ""
    elif block_type in {"image", "file", "pdf", "video"}:
        file_info = payload.get("external") or payload.get("file") or {}
        url = file_info.get("url")
        label = text or block_type
        return f"[{label}]({url})" if url else label
    rendered = f"{prefix}{text}".strip()
    child_markdown = ""
    if block.get("has_children"):
        children = await _notion_block_children(session, connection, str(block["id"]))
        child_parts = [await _notion_block_to_markdown(session, connection, child, depth=depth + 1) for child in children]
        child_markdown = "\n\n".join(part for part in child_parts if part.strip())
    if rendered and child_markdown:
        return f"{rendered}\n\n{child_markdown}".strip()
    return rendered or child_markdown


async def _notion_page_markdown(session: AsyncSession, connection: ConnectorConnection, page: dict[str, Any]) -> str:
    title = _notion_title_from_page(page)
    property_lines: list[str] = []
    for key, value in (page.get("properties") or {}).items():
        if not isinstance(value, dict):
            continue
        value_type = value.get("type")
        if value_type == "title":
            continue
        rendered = ""
        if value_type == "rich_text":
            rendered = _notion_rich_text(value.get("rich_text"))
        elif value_type == "select":
            selected = value.get("select") or {}
            rendered = str(selected.get("name") or "")
        elif value_type == "multi_select":
            rendered = ", ".join(str(item.get("name") or "") for item in value.get("multi_select") or [])
        elif value_type == "status":
            status_item = value.get("status") or {}
            rendered = str(status_item.get("name") or "")
        elif value_type == "date":
            date_item = value.get("date") or {}
            rendered = str(date_item.get("start") or "")
        elif value_type == "checkbox":
            rendered = "true" if value.get("checkbox") else "false"
        elif value_type == "number":
            rendered = "" if value.get("number") is None else str(value.get("number"))
        elif value_type == "url":
            rendered = str(value.get("url") or "")
        elif value_type == "email":
            rendered = str(value.get("email") or "")
        elif value_type == "phone_number":
            rendered = str(value.get("phone_number") or "")
        elif value_type == "people":
            rendered = ", ".join(str(item.get("name") or item.get("id") or "") for item in value.get("people") or [])
        if rendered:
            property_lines.append(f"- {key}: {rendered}")

    blocks = await _notion_block_children(session, connection, str(page["id"]))
    content_parts = [await _notion_block_to_markdown(session, connection, block) for block in blocks]
    content = "\n\n".join(part for part in content_parts if part.strip())
    sections = [f"# {title}"]
    if property_lines:
        sections.append("## 속성\n" + "\n".join(property_lines))
    if content:
        sections.append(content)
    return "\n\n".join(section for section in sections if section.strip()).strip()


async def _google_oauth_start(session: AsyncSession, *, state: str, code_verifier: str) -> dict[str, str]:
    params = httpx.QueryParams(
        {
            "client_id": get_settings().google_oauth_client_id,
            "redirect_uri": _app_callback_path("/api/connectors/google-drive/oauth/callback"),
            "response_type": "code",
            "scope": " ".join(GOOGLE_CONNECTOR_SCOPES),
            "state": state,
            "code_challenge": create_code_challenge(code_verifier),
            "code_challenge_method": "S256",
            "access_type": "offline",
            "prompt": "consent",
        }
    )
    return {"authorization_url": f"{GOOGLE_AUTH_URL}?{params}", "state": state}


async def _google_oauth_complete(*, code: str, code_verifier: str) -> ProviderOAuthResult:
    token_data = await _exchange_google_code(
        code=code,
        code_verifier=code_verifier,
        redirect_uri=_app_callback_path("/api/connectors/google-drive/oauth/callback"),
    )
    userinfo = await _google_userinfo(str(token_data["access_token"]))
    return ProviderOAuthResult(
        account_subject=str(userinfo["sub"]),
        display_name=f"Google Drive ({userinfo.get('email') or 'account'})",
        account_email=str(userinfo.get("email") or ""),
        encrypted_access_token=encrypt_secret(str(token_data["access_token"])) or "",
        encrypted_refresh_token=encrypt_secret(str(token_data.get("refresh_token") or "")),
        token_expires_at=future_utc(seconds=int(token_data.get("expires_in", 3600))),
        granted_scopes=list(GOOGLE_CONNECTOR_SCOPES),
    )


async def _github_oauth_start(session: AsyncSession, *, state: str, code_verifier: str) -> dict[str, str]:
    del session, code_verifier
    settings = get_settings()
    redirect_uri = settings.github_oauth_redirect_uri or _app_callback_path("/api/connectors/github/oauth/callback")
    params = httpx.QueryParams(
        {
            "client_id": settings.github_oauth_client_id,
            "redirect_uri": redirect_uri,
            "scope": " ".join(GITHUB_CONNECTOR_SCOPES),
            "state": state,
        }
    )
    return {"authorization_url": f"{GITHUB_AUTH_URL}?{params}", "state": state}


async def _github_oauth_complete(*, code: str, code_verifier: str) -> ProviderOAuthResult:
    del code_verifier
    settings = get_settings()
    redirect_uri = settings.github_oauth_redirect_uri or _app_callback_path("/api/connectors/github/oauth/callback")
    token_data = await _exchange_github_code(code=code, redirect_uri=redirect_uri)
    access_token = str(token_data["access_token"])
    userinfo = await _github_user(access_token)
    account_email = userinfo.get("email")
    if not account_email:
        account_email = await _github_primary_email(access_token)
    granted_scopes = [
        scope
        for scope in str(token_data.get("scope") or "").replace(",", " ").split()
        if scope
    ] or list(GITHUB_CONNECTOR_SCOPES)
    login = str(userinfo.get("login") or userinfo.get("name") or "account")
    return ProviderOAuthResult(
        account_subject=str(userinfo.get("id") or login),
        display_name=f"GitHub ({login})",
        account_email=str(account_email) if account_email else None,
        encrypted_access_token=encrypt_secret(access_token) or "",
        encrypted_refresh_token=None,
        token_expires_at=None,
        granted_scopes=granted_scopes,
    )


async def _notion_oauth_start(session: AsyncSession, *, state: str, code_verifier: str) -> dict[str, str]:
    del session, code_verifier
    params = httpx.QueryParams(
        {
            "owner": "user",
            "client_id": get_settings().notion_oauth_client_id,
            "redirect_uri": _app_callback_path("/api/connectors/notion/oauth/callback"),
            "response_type": "code",
            "state": state,
        }
    )
    return {"authorization_url": f"{NOTION_AUTH_URL}?{params}", "state": state}


async def _notion_oauth_complete(*, code: str, code_verifier: str) -> ProviderOAuthResult:
    del code_verifier
    token_data = await _exchange_notion_code(
        code=code,
        redirect_uri=_app_callback_path("/api/connectors/notion/oauth/callback"),
    )
    workspace_name = str(token_data.get("workspace_name") or "workspace")
    workspace_id = str(token_data.get("workspace_id") or token_data.get("bot_id") or workspace_name)
    return ProviderOAuthResult(
        account_subject=workspace_id,
        display_name=f"Notion ({workspace_name})",
        account_email=None,
        encrypted_access_token=encrypt_secret(str(token_data["access_token"])) or "",
        encrypted_refresh_token=None,
        token_expires_at=None,
        granted_scopes=[],
    )


async def _start_oauth_for_provider(session: AsyncSession, *, provider: str, state: str, code_verifier: str) -> dict[str, str]:
    if provider == ConnectorProvider.google_drive.value:
        return await _google_oauth_start(session, state=state, code_verifier=code_verifier)
    if provider == ConnectorProvider.github.value:
        return await _github_oauth_start(session, state=state, code_verifier=code_verifier)
    if provider == ConnectorProvider.notion.value:
        return await _notion_oauth_start(session, state=state, code_verifier=code_verifier)
    raise ConnectorError("Unsupported connector provider.")


async def _complete_oauth_for_provider(*, provider: str, code: str, code_verifier: str) -> ProviderOAuthResult:
    if provider == ConnectorProvider.google_drive.value:
        return await _google_oauth_complete(code=code, code_verifier=code_verifier)
    if provider == ConnectorProvider.github.value:
        return await _github_oauth_complete(code=code, code_verifier=code_verifier)
    if provider == ConnectorProvider.notion.value:
        return await _notion_oauth_complete(code=code, code_verifier=code_verifier)
    raise ConnectorError("Unsupported connector provider.")


async def _list_google_drive_files_paginated(
    session: AsyncSession,
    connection: ConnectorConnection,
    *,
    q: str,
    drive_id: str | None = None,
) -> list[dict[str, Any]]:
    files: list[dict[str, Any]] = []
    page_token: str | None = None
    while True:
        params: dict[str, Any] = {
            "q": q,
            "pageSize": 1000,
            "fields": "nextPageToken,files(id,name,mimeType,modifiedTime,version,md5Checksum,webViewLink,fileExtension,parents,trashed,driveId)",
            "supportsAllDrives": "true",
            "includeItemsFromAllDrives": "true",
        }
        if drive_id:
            params["corpora"] = "drive"
            params["driveId"] = drive_id
        if page_token:
            params["pageToken"] = page_token
        data = await _google_json(session, connection, "/files", params=params)
        files.extend(list(data.get("files", [])))
        page_token = data.get("nextPageToken")
        if not page_token:
            break
    return files


def _markdown_table(rows: list[list[Any]]) -> str:
    if not rows:
        return ""
    normalized = [["" if cell is None else str(cell).replace("\n", " ").strip() for cell in row] for row in rows]
    header = normalized[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in normalized[1:]:
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[: len(header)]) + " |")
    return "\n".join(lines)


def _extract_xlsx(data: bytes) -> ExtractedContent:
    workbook = load_workbook(io.BytesIO(data), data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"## {sheet.title}")
        rows = [[cell for cell in row] for row in sheet.iter_rows(values_only=True)]
        parts.append(_markdown_table(rows) if rows else "_빈 시트_")
    return ExtractedContent(content_type="markdown", content="\n\n".join(parts).strip(), doc_type="data")


def _extract_docx(data: bytes) -> ExtractedContent:
    document = DocxDocument(io.BytesIO(data))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return ExtractedContent(content_type="text", content="\n\n".join(paragraphs).strip(), doc_type="knowledge")


def _extract_pptx(data: bytes) -> ExtractedContent:
    presentation = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = [f"## Slide {index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and str(shape.text).strip():
                lines.append(str(shape.text).strip())
        slides.append("\n\n".join(lines))
    return ExtractedContent(content_type="markdown", content="\n\n".join(slides).strip(), doc_type="knowledge")


def _extract_pdf(data: bytes) -> ExtractedContent:
    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return ExtractedContent(content_type="text", content=normalize_whitespace("\n\n".join(pages)), doc_type="knowledge")


def _extract_image(data: bytes) -> ExtractedContent:
    image = Image.open(io.BytesIO(data))
    text = pytesseract.image_to_string(image)
    return ExtractedContent(content_type="text", content=normalize_whitespace(text), doc_type="knowledge")


def _extract_csv(data: bytes) -> ExtractedContent:
    decoded = data.decode("utf-8-sig", errors="ignore")
    rows = list(csv.reader(io.StringIO(decoded)))
    return ExtractedContent(content_type="markdown", content=_markdown_table(rows), doc_type="data")


def extract_file_content(file_meta: dict[str, Any], data: bytes) -> ExtractedContent:
    mime_type = str(file_meta.get("mimeType") or "")
    name = str(file_meta.get("name") or "document")
    lower_name = name.lower()

    if mime_type == "text/html" or lower_name.endswith((".html", ".htm")):
        return ExtractedContent(content_type="html", content=data.decode("utf-8", errors="ignore"), doc_type="knowledge")
    if mime_type in {"text/plain", "text/markdown"} or lower_name.endswith((".txt", ".md", ".markdown")):
        doc_type = "data" if lower_name.endswith(".csv") else "knowledge"
        content_type = "markdown" if lower_name.endswith((".md", ".markdown")) else "text"
        return ExtractedContent(content_type=content_type, content=data.decode("utf-8-sig", errors="ignore"), doc_type=doc_type)
    if mime_type == "text/csv" or lower_name.endswith(".csv"):
        return _extract_csv(data)
    if mime_type == "application/pdf" or lower_name.endswith(".pdf"):
        return _extract_pdf(data)
    if lower_name.endswith(".docx"):
        return _extract_docx(data)
    if lower_name.endswith(".xlsx"):
        return _extract_xlsx(data)
    if lower_name.endswith(".pptx"):
        return _extract_pptx(data)
    if mime_type.startswith("image/") or lower_name.endswith((".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff")):
        return _extract_image(data)
    raise ConnectorError(f"Unsupported file type: {mime_type or lower_name}")


async def _download_google_file(session: AsyncSession, connection: ConnectorConnection, file_meta: dict[str, Any]) -> ExtractedContent:
    file_id = str(file_meta["id"])
    mime_type = str(file_meta.get("mimeType") or "")
    if mime_type == GOOGLE_DOC_MIME:
        data = await _google_bytes(session, connection, f"/files/{file_id}/export", params={"mimeType": "text/html"})
        return ExtractedContent(content_type="html", content=data.decode("utf-8", errors="ignore"), doc_type="knowledge")
    if mime_type == GOOGLE_SHEET_MIME:
        data = await _google_bytes(session, connection, f"/files/{file_id}/export", params={"mimeType": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"})
        return _extract_xlsx(data)
    if mime_type == GOOGLE_SLIDE_MIME:
        data = await _google_bytes(session, connection, f"/files/{file_id}/export", params={"mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation"})
        return _extract_pptx(data)
    if mime_type == GOOGLE_DRAWING_MIME:
        data = await _google_bytes(session, connection, f"/files/{file_id}/export", params={"mimeType": "application/pdf"})
        return _extract_pdf(data)
    data = await _google_bytes(session, connection, f"/files/{file_id}", params={"alt": "media", "supportsAllDrives": "true"})
    return extract_file_content(file_meta, data)


async def _google_browse(
    session: AsyncSession,
    connection: ConnectorConnection,
    *,
    kind: str,
    parent_id: str | None,
    container_id: str | None,
) -> ConnectorBrowseResponse:
    if kind == ConnectorResourceKind.shared_drive.value:
        data = await _google_json(session, connection, "/drives", params={"pageSize": 100})
        return ConnectorBrowseResponse(
            kind=kind,
            items=[
                ConnectorBrowseItem(
                    id=str(item["id"]),
                    name=str(item["name"]),
                    resource_kind=ConnectorResourceKind.shared_drive.value,
                    has_children=True,
                    provider_metadata={"drive_id": str(item["id"])},
                )
                for item in data.get("drives", [])
            ],
        )

    folder_parent = parent_id or "root"
    query = f"mimeType = '{GOOGLE_FOLDER_MIME}' and trashed = false and '{folder_parent}' in parents"
    params: dict[str, Any] = {
        "q": query,
        "pageSize": 100,
        "fields": "files(id,name,mimeType,parents,driveId,webViewLink)",
        "supportsAllDrives": "true",
        "includeItemsFromAllDrives": "true",
        "orderBy": "name",
    }
    if container_id:
        params["corpora"] = "drive"
        params["driveId"] = container_id
    data = await _google_json(session, connection, "/files", params=params)
    return ConnectorBrowseResponse(
        kind=kind,
        parent_external_id=parent_id,
        items=[
            ConnectorBrowseItem(
                id=str(item["id"]),
                name=str(item["name"]),
                resource_kind=ConnectorResourceKind.folder.value,
                resource_url=item.get("webViewLink"),
                parent_external_id=(item.get("parents") or [None])[0],
                has_children=True,
                provider_metadata={"drive_id": item.get("driveId")},
            )
            for item in data.get("files", [])
        ],
    )


async def _notion_search(
    session: AsyncSession,
    connection: ConnectorConnection,
    *,
    kind: str,
    query: str | None,
    cursor: str | None,
) -> ConnectorBrowseResponse:
    filter_value = "database" if kind == ConnectorResourceKind.database.value else "page"
    body: dict[str, Any] = {
        "page_size": 25,
        "filter": {
            "property": "object",
            "value": filter_value,
        },
        "sort": {
            "direction": "descending",
            "timestamp": "last_edited_time",
        },
    }
    if query:
        body["query"] = query
    if cursor:
        body["start_cursor"] = cursor
    payload = await _notion_request(session, connection, "POST", "/search", payload=body)
    items: list[ConnectorBrowseItem] = []
    for item in payload.get("results", []):
        resource_kind = ConnectorResourceKind.database.value if item.get("object") == "database" else ConnectorResourceKind.page.value
        name = _notion_title_from_database(item) if resource_kind == ConnectorResourceKind.database.value else _notion_title_from_page(item)
        items.append(
            ConnectorBrowseItem(
                id=str(item["id"]),
                name=name,
                resource_kind=resource_kind,
                resource_url=item.get("url"),
                has_children=resource_kind == ConnectorResourceKind.database.value,
                provider_metadata={"object": item.get("object")},
            )
        )
    return ConnectorBrowseResponse(
        kind=kind,
        cursor=payload.get("next_cursor"),
        has_more=bool(payload.get("has_more")),
        items=items,
    )


def _github_doc_path_supported(path: str) -> bool:
    normalized = path.strip("/")
    lower = normalized.lower()
    if not normalized:
        return False
    name = lower.rsplit("/", 1)[-1]
    if "/" not in normalized and name.startswith("readme"):
        return "." not in name or lower.endswith(tuple(GITHUB_DOC_EXTENSIONS))
    if not (lower.startswith("docs/") or lower.startswith("doc/")):
        return False
    return any(lower.endswith(extension) for extension in GITHUB_DOC_EXTENSIONS)


def _github_evidence_path_excluded(path: str) -> bool:
    lower = path.lower()
    excluded_segments = {
        ".git",
        ".github",
        ".next",
        ".venv",
        "__pycache__",
        "build",
        "coverage",
        "dist",
        "node_modules",
        "target",
        "tmp",
        "vendor",
    }
    parts = {part for part in PurePosixPath(lower).parts if part}
    return bool(parts & excluded_segments)


def _github_probably_binary(raw_content: bytes) -> bool:
    if not raw_content:
        return False
    if b"\x00" in raw_content[:2048]:
        return True
    try:
        raw_content.decode("utf-8-sig")
    except UnicodeDecodeError:
        return True
    return False


def _github_document_content_type(path: str) -> str:
    lower = path.lower()
    if lower.endswith((".md", ".mdx")) or lower.rsplit("/", 1)[-1].startswith("readme"):
        return "markdown"
    return "text"


def _github_document_title(repo_name: str, path: str) -> str:
    return f"{repo_name} · {path}"


def _github_evidence_title(repo_name: str, path: str) -> str:
    return f"{repo_name} · 근거 · {path}"


def _notion_export_title(path: str) -> str:
    normalized = path.replace("\\", "/").strip("/")
    if not normalized:
        return "Notion export"
    return normalized.rsplit("/", 1)[-1].rsplit(".", 1)[0] or "Notion export"


def _notion_export_supported_path(path: str) -> bool:
    normalized = path.replace("\\", "/").strip("/")
    if not normalized or normalized.startswith("__macosx/"):
        return False
    suffix = PurePosixPath(normalized).suffix.lower()
    return suffix in NOTION_EXPORT_EXTENSIONS


def _notion_export_content_type(path: str) -> str:
    suffix = PurePosixPath(path).suffix.lower()
    if suffix in {".md", ".markdown"}:
        return "markdown"
    if suffix in {".html", ".htm"}:
        return "html"
    return "text"


def _iter_notion_export_items(filename: str, raw_bytes: bytes) -> list[PreparedSyncItem]:
    items: list[PreparedSyncItem] = []
    if filename.lower().endswith(".zip"):
        with zipfile.ZipFile(io.BytesIO(raw_bytes)) as archive:
            for member in archive.infolist():
                if member.is_dir():
                    continue
                if not _notion_export_supported_path(member.filename):
                    continue
                content = archive.read(member.filename).decode("utf-8-sig", errors="ignore")
                items.append(
                    PreparedSyncItem(
                        external_item_id=member.filename,
                        title=_notion_export_title(member.filename),
                        source_url=None,
                        source_revision_id=None,
                        mime_type="application/octet-stream",
                        content_type=_notion_export_content_type(member.filename),
                        content=content,
                        doc_type="knowledge",
                        provider_metadata={
                            "notion_export_path": member.filename,
                            "corpus_role": "glossary_evidence",
                        },
                    )
                )
        return items

    content = raw_bytes.decode("utf-8-sig", errors="ignore")
    if not content.strip():
        return items
    items.append(
        PreparedSyncItem(
            external_item_id=filename,
            title=_notion_export_title(filename),
            source_url=None,
            source_revision_id=None,
            mime_type="application/octet-stream",
            content_type=_notion_export_content_type(filename),
            content=content,
            doc_type="knowledge",
            provider_metadata={
                "notion_export_path": filename,
                "corpus_role": "glossary_evidence",
            },
        )
    )
    return items


def _github_has_next_page(headers: httpx.Headers) -> bool:
    link = headers.get("Link") or headers.get("link") or ""
    return 'rel="next"' in link


async def _github_search_repositories(
    session: AsyncSession,
    connection: ConnectorConnection,
    *,
    kind: str,
    query: str | None,
    cursor: str | None,
) -> ConnectorBrowseResponse:
    try:
        page = max(int(cursor or "1"), 1)
    except ValueError:
        page = 1
    payload, headers = await _github_request(
        session,
        connection,
        "GET",
        "/user/repos",
        params={
            "per_page": 50,
            "page": page,
            "sort": "updated",
            "affiliation": "owner,collaborator,organization_member",
        },
    )
    repositories = payload if isinstance(payload, list) else []
    normalized_query = (query or "").strip().lower()
    items: list[ConnectorBrowseItem] = []
    for repository in repositories:
        if not isinstance(repository, dict):
            continue
        full_name = str(repository.get("full_name") or "")
        owner = str((repository.get("owner") or {}).get("login") or "")
        repo = str(repository.get("name") or "")
        if not owner or not repo or not full_name:
            continue
        haystack = " ".join(
            filter(
                None,
                [
                    full_name.lower(),
                    repo.lower(),
                    str(repository.get("description") or "").lower(),
                ],
            )
        )
        if normalized_query and normalized_query not in haystack:
            continue
        items.append(
            ConnectorBrowseItem(
                id=full_name,
                name=full_name,
                resource_kind=kind,
                resource_url=repository.get("html_url"),
                has_children=False,
                provider_metadata={
                    "owner": owner,
                    "repo": repo,
                    "default_branch": repository.get("default_branch"),
                    "archived": bool(repository.get("archived")),
                },
            )
        )
    next_cursor = str(page + 1) if _github_has_next_page(headers) else None
    return ConnectorBrowseResponse(
        kind=kind,
        items=items,
        cursor=next_cursor,
        has_more=next_cursor is not None,
    )


async def _github_sync_items(
    session: AsyncSession,
    connection: ConnectorConnection,
    resource: ConnectorResource,
) -> list[PreparedSyncItem]:
    if resource.resource_kind not in {
        ConnectorResourceKind.repository_docs.value,
        ConnectorResourceKind.repository_evidence.value,
    }:
        raise ConnectorError("Unsupported GitHub resource kind.")
    metadata = dict(resource.provider_metadata or {})
    owner = str(metadata.get("owner") or "")
    repo = str(metadata.get("repo") or "")
    if (not owner or not repo) and "/" in resource.external_id:
        owner, repo = resource.external_id.split("/", 1)
    if not owner or not repo:
        raise ConnectorError("GitHub repository identifier is invalid.")

    repository, _headers = await _github_request(session, connection, "GET", f"/repos/{owner}/{repo}")
    default_branch = str(metadata.get("default_branch") or repository.get("default_branch") or "main")
    repo_name = str(repository.get("name") or repo)
    repo_html_url = str(repository.get("html_url") or f"https://github.com/{owner}/{repo}")

    tree_ref = quote(default_branch, safe="")
    tree_payload, _tree_headers = await _github_request(
        session,
        connection,
        "GET",
        f"/repos/{owner}/{repo}/git/trees/{tree_ref}",
        params={"recursive": "1"},
    )
    tree_items = tree_payload.get("tree") or []
    if resource.resource_kind == ConnectorResourceKind.repository_docs.value:
        candidate_paths = sorted(
            str(item.get("path"))
            for item in tree_items
            if isinstance(item, dict)
            and item.get("type") == "blob"
            and _github_doc_path_supported(str(item.get("path") or ""))
        )
    else:
        candidate_paths = sorted(
            str(item.get("path"))
            for item in tree_items
            if isinstance(item, dict)
            and item.get("type") == "blob"
            and not _github_evidence_path_excluded(str(item.get("path") or ""))
        )

    prepared: list[PreparedSyncItem] = []
    for path in candidate_paths:
        encoded_path = quote(path, safe="/")
        file_payload, _file_headers = await _github_request(
            session,
            connection,
            "GET",
            f"/repos/{owner}/{repo}/contents/{encoded_path}",
            params={"ref": default_branch},
        )
        if isinstance(file_payload, list):
            continue
        raw_content: bytes
        if file_payload.get("content") and file_payload.get("encoding") == "base64":
            raw_content = base64.b64decode(str(file_payload["content"]).encode("utf-8"))
        elif file_payload.get("download_url"):
            raw_content = await _github_raw_bytes(session, connection, str(file_payload["download_url"]))
        else:
            raise ConnectorError(f"GitHub file content is unavailable for {path}.")
        if resource.resource_kind == ConnectorResourceKind.repository_evidence.value and _github_probably_binary(raw_content):
            continue
        content = raw_content.decode("utf-8-sig", errors="ignore")
        canonical_url = str(file_payload.get("html_url") or f"{repo_html_url}/blob/{default_branch}/{path}")
        content_type = _github_document_content_type(path)
        prepared.append(
            PreparedSyncItem(
                external_item_id=path,
                title=(
                    _github_document_title(repo_name, path)
                    if resource.resource_kind == ConnectorResourceKind.repository_docs.value
                    else _github_evidence_title(repo_name, path)
                ),
                source_url=canonical_url,
                source_revision_id=str(file_payload.get("sha") or ""),
                mime_type="text/markdown" if content_type == "markdown" else "text/plain",
                content_type="markdown" if path.lower().endswith((".md", ".mdx")) else "text",
                content=content,
                doc_type="knowledge",
                provider_metadata={
                    "github_owner": owner,
                    "github_repo": repo,
                    "github_path": path,
                    "github_default_branch": default_branch,
                    "github_html_url": canonical_url,
                    "corpus_role": (
                        "docs"
                        if resource.resource_kind == ConnectorResourceKind.repository_docs.value
                        else "glossary_evidence"
                    ),
                },
            )
        )
    return prepared


async def _google_sync_items(
    session: AsyncSession,
    connection: ConnectorConnection,
    resource: ConnectorResource,
) -> list[PreparedSyncItem]:
    if resource.resource_kind == ConnectorResourceKind.shared_drive.value:
        files = [
            item
            for item in await _list_google_drive_files_paginated(
                session,
                connection,
                q="trashed = false",
                drive_id=resource.external_id,
            )
            if item.get("mimeType") != GOOGLE_FOLDER_MIME
        ]
    else:
        queue: deque[str] = deque([resource.external_id])
        seen_folders: set[str] = set()
        files: list[dict[str, Any]] = []
        while queue:
            folder_id = queue.popleft()
            if folder_id in seen_folders:
                continue
            seen_folders.add(folder_id)
            children = await _list_google_drive_files_paginated(
                session,
                connection,
                q=f"trashed = false and '{folder_id}' in parents",
            )
            for child in children:
                if child.get("mimeType") == GOOGLE_FOLDER_MIME:
                    if resource.sync_children:
                        queue.append(str(child["id"]))
                    continue
                files.append(child)

    prepared: list[PreparedSyncItem] = []
    for file_meta in files:
        try:
            extracted = await _download_google_file(session, connection, file_meta)
            prepared.append(
                PreparedSyncItem(
                    external_item_id=str(file_meta["id"]),
                    title=str(file_meta.get("name") or file_meta["id"]),
                    source_url=file_meta.get("webViewLink"),
                    source_revision_id=_google_revision_token(file_meta),
                    mime_type=file_meta.get("mimeType"),
                    content_type=extracted.content_type,
                    content=extracted.content,
                    doc_type=extracted.doc_type,
                    provider_metadata={
                        "google_mime_type": file_meta.get("mimeType"),
                        "drive_id": file_meta.get("driveId"),
                    },
                )
            )
        except ConnectorError as exc:
            prepared.append(
                PreparedSyncItem(
                    external_item_id=str(file_meta["id"]),
                    title=str(file_meta.get("name") or file_meta["id"]),
                    source_url=file_meta.get("webViewLink"),
                    source_revision_id=_google_revision_token(file_meta),
                    mime_type=file_meta.get("mimeType"),
                    content_type=None,
                    content=None,
                    doc_type="knowledge",
                    unsupported_reason=str(exc),
                    provider_metadata={"google_mime_type": file_meta.get("mimeType")},
                )
            )
    return prepared


async def _notion_sync_page_item(
    session: AsyncSession,
    connection: ConnectorConnection,
    page: dict[str, Any],
    *,
    database_id: str | None = None,
) -> PreparedSyncItem:
    markdown = await _notion_page_markdown(session, connection, page)
    return PreparedSyncItem(
        external_item_id=str(page["id"]),
        title=_notion_title_from_page(page),
        source_url=page.get("url"),
        source_revision_id=_notion_revision_token(page),
        mime_type="notion/page",
        content_type="markdown",
        content=markdown,
        doc_type="knowledge",
        provider_metadata={
            "object": "page",
            "database_id": database_id,
        },
    )


async def _notion_sync_items(
    session: AsyncSession,
    connection: ConnectorConnection,
    resource: ConnectorResource,
) -> list[PreparedSyncItem]:
    if resource.resource_kind == ConnectorResourceKind.page.value:
        page = await _notion_request(session, connection, "GET", f"/pages/{resource.external_id}")
        return [await _notion_sync_page_item(session, connection, page)]

    database = await _notion_request(session, connection, "GET", f"/databases/{resource.external_id}")
    items: list[PreparedSyncItem] = []
    cursor: str | None = None
    while True:
        body: dict[str, Any] = {"page_size": 100}
        if cursor:
            body["start_cursor"] = cursor
        payload = await _notion_request(session, connection, "POST", f"/databases/{resource.external_id}/query", payload=body)
        for page in payload.get("results", []):
            items.append(await _notion_sync_page_item(session, connection, page, database_id=str(database["id"])))
        if not payload.get("has_more"):
            break
        cursor = payload.get("next_cursor")
        if not cursor:
            break
    return items


async def _browse_for_provider(
    session: AsyncSession,
    connection: ConnectorConnection,
    *,
    kind: str,
    query: str | None,
    cursor: str | None,
    parent_id: str | None,
    container_id: str | None,
) -> ConnectorBrowseResponse:
    if connection.provider == ConnectorProvider.google_drive.value:
        return await _google_browse(session, connection, kind=kind, parent_id=parent_id, container_id=container_id)
    if connection.provider == ConnectorProvider.github.value:
        return await _github_search_repositories(session, connection, kind=kind, query=query, cursor=cursor)
    if connection.provider == ConnectorProvider.notion.value:
        return await _notion_search(session, connection, kind=kind, query=query, cursor=cursor)
    raise ConnectorError("Unsupported connector provider.")


async def _sync_items_for_resource(
    session: AsyncSession,
    connection: ConnectorConnection,
    resource: ConnectorResource,
) -> list[PreparedSyncItem]:
    if connection.provider == ConnectorProvider.google_drive.value:
        return await _google_sync_items(session, connection, resource)
    if connection.provider == ConnectorProvider.github.value:
        return await _github_sync_items(session, connection, resource)
    if connection.provider == ConnectorProvider.notion.value:
        return await _notion_sync_items(session, connection, resource)
    raise ConnectorError("Unsupported connector provider.")


async def _get_connection_or_raise(
    session: AsyncSession,
    connection_id: UUID,
    auth_user: AuthenticatedUser,
    *,
    allow_workspace_read: bool = True,
) -> ConnectorConnection:
    connection = await session.get(ConnectorConnection, connection_id)
    if connection is None:
        raise ConnectorNotFoundError("Connector not found.")
    if auth_user.current_workspace_id is None or connection.workspace_id != auth_user.current_workspace_id:
        raise ConnectorForbiddenError("Connector does not belong to the current workspace.")
    if connection.owner_scope == ConnectorOwnerScope.workspace.value:
        if not allow_workspace_read:
            raise ConnectorForbiddenError("Workspace connector access is not allowed.")
        return connection
    if connection.owner_user_id != auth_user.user.id:
        raise ConnectorForbiddenError("Personal connector access denied.")
    return connection


async def _get_resource_or_raise(session: AsyncSession, resource_id: UUID) -> ConnectorResource:
    resource = await session.get(ConnectorResource, resource_id)
    if resource is None:
        raise ConnectorNotFoundError("Connector resource not found.")
    return resource


def _ensure_scope_permission(scope: str, auth_user: AuthenticatedUser) -> None:
    _validate_owner_scope(scope)
    if scope == ConnectorOwnerScope.workspace.value and not auth_user.can_manage_workspace_connectors:
        raise ConnectorForbiddenError("Workspace connectors can only be managed by workspace owners or admins.")


async def get_connectors_readiness(
    session: AsyncSession,
    auth_user: AuthenticatedUser | None,
) -> ConnectorReadinessResponse:
    items: list[ConnectorProviderReadiness] = []
    workspace_id = auth_user.current_workspace_id if auth_user is not None else None
    for provider in [ConnectorProvider.google_drive.value, ConnectorProvider.github.value, ConnectorProvider.notion.value]:
        workspace_connection: ConnectorConnection | None = None
        healthy_source_count = 0
        needs_attention_count = 0
        if workspace_id is not None:
            workspace_connection = (
                await session.execute(
                    select(ConnectorConnection)
                    .where(
                        ConnectorConnection.provider == provider,
                        ConnectorConnection.workspace_id == workspace_id,
                        ConnectorConnection.owner_scope == ConnectorOwnerScope.workspace.value,
                    )
                    .order_by(ConnectorConnection.created_at.asc())
                    .limit(1)
                )
            ).scalar_one_or_none()
            if workspace_connection is not None:
                resources = list(
                    (
                        await session.execute(
                            select(ConnectorResource).where(ConnectorResource.connection_id == workspace_connection.id)
                        )
                    ).scalars().all()
                )
                healthy_source_count = sum(
                    1
                    for resource in resources
                    if resource.status == ConnectorResourceStatus.active.value
                    and int((resource.last_sync_summary or {}).get("failed", 0)) == 0
                )
                needs_attention_count = sum(
                    1
                    for resource in resources
                    if resource.status != ConnectorResourceStatus.active.value
                    or int((resource.last_sync_summary or {}).get("failed", 0)) > 0
                )
        if workspace_connection is not None and workspace_connection.status != ConnectorStatus.active.value:
            needs_attention_count += 1
        items.append(
            _provider_readiness_summary(
                provider,
                connection=workspace_connection,
                auth_user=auth_user,
                healthy_source_count=healthy_source_count,
                needs_attention_count=needs_attention_count,
            )
        )
    return ConnectorReadinessResponse(providers=items)


async def list_connections(session: AsyncSession, auth_user: AuthenticatedUser, *, scope: str) -> ConnectorListResponse:
    scope = _validate_owner_scope(scope)
    if auth_user.current_workspace_id is None:
        return ConnectorListResponse(items=[])
    if scope == ConnectorOwnerScope.workspace.value:
        stmt = select(ConnectorConnection).where(
            ConnectorConnection.workspace_id == auth_user.current_workspace_id,
            ConnectorConnection.owner_scope == ConnectorOwnerScope.workspace.value,
        )
    else:
        stmt = select(ConnectorConnection).where(
            ConnectorConnection.workspace_id == auth_user.current_workspace_id,
            ConnectorConnection.owner_scope == ConnectorOwnerScope.personal.value,
            ConnectorConnection.owner_user_id == auth_user.user.id,
        )
    connections = list((await session.execute(stmt.order_by(ConnectorConnection.provider.asc(), ConnectorConnection.created_at.asc()))).scalars().all())
    connection_ids = [connection.id for connection in connections]
    resources = list(
        (
            await session.execute(
                select(ConnectorResource)
                .where(ConnectorResource.connection_id.in_(connection_ids))
                .order_by(ConnectorResource.provider.asc(), ConnectorResource.name.asc())
            )
        ).scalars().all()
    ) if connection_ids else []
    resources_by_connection: dict[UUID, list[ConnectorResource]] = {}
    for resource in resources:
        resources_by_connection.setdefault(resource.connection_id, []).append(resource)
    return ConnectorListResponse(
        items=[
            _connection_summary(connection, resources_by_connection.get(connection.id, []))
            for connection in connections
        ]
    )


async def get_connection_detail(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    connection_id: UUID,
) -> ConnectorConnectionSummary:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    resources = list(
        (
            await session.execute(
                select(ConnectorResource)
                .where(ConnectorResource.connection_id == connection.id)
                .order_by(ConnectorResource.provider.asc(), ConnectorResource.name.asc())
            )
        ).scalars().all()
    )
    return _connection_summary(connection, resources)


async def update_connection(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    connection_id: UUID,
    payload: ConnectorUpdateRequest,
) -> ConnectorConnectionSummary:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    if connection.owner_scope == ConnectorOwnerScope.workspace.value:
        _ensure_scope_permission(connection.owner_scope, auth_user)
    if payload.display_name is not None:
        connection.display_name = payload.display_name.strip() or connection.display_name
    if payload.status is not None:
        connection.status = payload.status
    await session.commit()
    return await get_connection_detail(session, auth_user, connection_id)


async def delete_connection(session: AsyncSession, auth_user: AuthenticatedUser, connection_id: UUID) -> None:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    if connection.owner_scope == ConnectorOwnerScope.workspace.value:
        _ensure_scope_permission(connection.owner_scope, auth_user)
    await session.delete(connection)
    await session.commit()


async def start_provider_oauth(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    *,
    provider: str,
    owner_scope: str,
    return_path: str = "/connectors",
) -> dict[str, str]:
    provider = _normalize_provider(provider)
    _ensure_provider_configured(provider)
    owner_scope = _validate_owner_scope(owner_scope)
    _ensure_scope_permission(owner_scope, auth_user)
    if auth_user.current_workspace_id is None:
        raise ConnectorForbiddenError("Workspace context is required.")
    verifier = generate_code_verifier()
    state = generate_state_token()
    session.add(
        ConnectorOAuthState(
            state=state,
            purpose=ConnectorOAuthPurpose.connect_provider.value,
            workspace_id=auth_user.current_workspace_id,
            owner_scope=owner_scope,
            owner_user_id=auth_user.user.id if owner_scope == ConnectorOwnerScope.personal.value else None,
            code_verifier=verifier,
            return_path=_safe_return_path(return_path),
            expires_at=future_utc(seconds=get_settings().oauth_state_ttl_seconds),
        )
    )
    await session.commit()
    return await _start_oauth_for_provider(session, provider=provider, state=state, code_verifier=verifier)


async def complete_provider_oauth(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    *,
    provider: str,
    state: str,
    code: str,
) -> ConnectorOAuthCallbackResponse:
    provider = _normalize_provider(provider)
    state_row = (
        await session.execute(
            select(ConnectorOAuthState).where(
                ConnectorOAuthState.state == state,
                ConnectorOAuthState.purpose == ConnectorOAuthPurpose.connect_provider.value,
            )
        )
    ).scalar_one_or_none()
    if state_row is None or state_row.expires_at < utcnow():
        raise ConnectorError("Connector OAuth state is invalid or expired.")
    if auth_user.current_workspace_id is None or state_row.workspace_id != auth_user.current_workspace_id:
        raise ConnectorForbiddenError("Connector callback workspace does not match the current workspace.")
    if state_row.owner_scope == ConnectorOwnerScope.workspace.value and not auth_user.can_manage_workspace_connectors:
        raise ConnectorForbiddenError("Workspace connector callback requires a workspace owner or admin.")
    if state_row.owner_scope == ConnectorOwnerScope.personal.value and state_row.owner_user_id != auth_user.user.id:
        raise ConnectorForbiddenError("Connector callback user does not match the original owner.")

    oauth_result = await _complete_oauth_for_provider(
        provider=provider,
        code=code,
        code_verifier=state_row.code_verifier,
    )
    owner_user_id = auth_user.user.id if state_row.owner_scope == ConnectorOwnerScope.personal.value else None
    existing = (
        await session.execute(
            select(ConnectorConnection).where(
                ConnectorConnection.provider == provider,
                ConnectorConnection.workspace_id == state_row.workspace_id,
                ConnectorConnection.owner_scope == state_row.owner_scope,
                ConnectorConnection.owner_user_id == owner_user_id,
            )
        )
    ).scalar_one_or_none()
    if existing is None:
        existing = ConnectorConnection(
            provider=provider,
            workspace_id=state_row.workspace_id,
            owner_scope=state_row.owner_scope,
            owner_user_id=owner_user_id,
            display_name=oauth_result.display_name,
            account_email=oauth_result.account_email,
            account_subject=oauth_result.account_subject,
            status=ConnectorStatus.active.value,
            encrypted_access_token=oauth_result.encrypted_access_token,
            encrypted_refresh_token=oauth_result.encrypted_refresh_token,
            token_expires_at=oauth_result.token_expires_at,
            granted_scopes=oauth_result.granted_scopes,
            last_validated_at=utcnow(),
        )
        session.add(existing)
        await session.flush()
    else:
        existing.display_name = oauth_result.display_name
        existing.account_email = oauth_result.account_email
        existing.account_subject = oauth_result.account_subject
        existing.status = ConnectorStatus.active.value
        existing.encrypted_access_token = oauth_result.encrypted_access_token
        existing.encrypted_refresh_token = oauth_result.encrypted_refresh_token
        existing.token_expires_at = oauth_result.token_expires_at
        existing.granted_scopes = oauth_result.granted_scopes
        existing.last_validated_at = utcnow()
        await session.flush()
    await session.delete(state_row)
    await session.commit()
    return ConnectorOAuthCallbackResponse(
        redirect_to=state_row.return_path or "/connectors",
        connection=await get_connection_detail(session, auth_user, existing.id),
    )


async def browse_connection(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    connection_id: UUID,
    *,
    kind: str | None,
    query: str | None = None,
    cursor: str | None = None,
    parent_id: str | None = None,
    container_id: str | None = None,
) -> ConnectorBrowseResponse:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    browse_kind = _validate_browse_kind(connection.provider, kind)
    return await _browse_for_provider(
        session,
        connection,
        kind=browse_kind,
        query=query,
        cursor=cursor,
        parent_id=parent_id,
        container_id=container_id,
    )


async def create_resource(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    connection_id: UUID,
    payload: ConnectorResourceCreateRequest,
) -> ConnectorResourceSummary:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    if connection.owner_scope == ConnectorOwnerScope.workspace.value:
        _ensure_scope_permission(connection.owner_scope, auth_user)
    resource_kind = _validate_resource_kind(connection.provider, payload.resource_kind)
    visibility_scope = _normalize_visibility_scope(connection.provider, resource_kind, payload.visibility_scope)
    selection_mode = _normalize_selection_mode(connection.provider, resource_kind, payload.selection_mode)
    sync_children, sync_mode, interval = _resource_sync_defaults(connection, payload)
    existing = (
        await session.execute(
            select(ConnectorResource).where(
                ConnectorResource.connection_id == connection.id,
                ConnectorResource.resource_kind == resource_kind,
                ConnectorResource.external_id == payload.external_id,
            )
        )
    ).scalar_one_or_none()
    resource = existing or ConnectorResource(
        connection_id=connection.id,
        provider=connection.provider,
        resource_kind=resource_kind,
        external_id=payload.external_id,
        name=payload.name,
        resource_url=payload.resource_url,
        parent_external_id=payload.parent_external_id,
        visibility_scope=visibility_scope,
        selection_mode=selection_mode,
        sync_children=sync_children,
        sync_mode=sync_mode,
        sync_interval_minutes=interval,
        status=ConnectorResourceStatus.active.value,
        next_auto_sync_at=future_utc(seconds=interval * 60) if sync_mode == ConnectorSyncMode.auto.value and interval else None,
        provider_metadata=payload.provider_metadata,
    )
    if existing is None:
        session.add(resource)
    else:
        resource.name = payload.name
        resource.resource_url = payload.resource_url
        resource.parent_external_id = payload.parent_external_id
        resource.visibility_scope = visibility_scope
        resource.selection_mode = selection_mode
        resource.sync_children = sync_children
        resource.sync_mode = sync_mode
        resource.sync_interval_minutes = interval
        resource.next_auto_sync_at = future_utc(seconds=interval * 60) if sync_mode == ConnectorSyncMode.auto.value and interval else None
        resource.status = ConnectorResourceStatus.active.value
        resource.provider_metadata = payload.provider_metadata
    await session.commit()
    await session.refresh(resource)
    return _resource_summary(resource)


async def import_notion_export_resource(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    connection_id: UUID,
    *,
    name: str,
    filename: str,
    content_bytes: bytes,
    visibility_scope: str | None = None,
) -> ConnectorResourceSummary:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    if connection.owner_scope == ConnectorOwnerScope.workspace.value:
        _ensure_scope_permission(connection.owner_scope, auth_user)
    if connection.provider != ConnectorProvider.notion.value:
        raise ConnectorError("Notion export upload is only supported for Notion connections.")

    resource_kind = ConnectorResourceKind.export_upload.value
    normalized_visibility = _normalize_visibility_scope(connection.provider, resource_kind, visibility_scope)
    resource = ConnectorResource(
        connection_id=connection.id,
        provider=connection.provider,
        resource_kind=resource_kind,
        external_id=f"export:{uuid4()}",
        name=name.strip() or _notion_export_title(filename),
        resource_url=None,
        parent_external_id=None,
        visibility_scope=normalized_visibility,
        selection_mode="export_upload",
        sync_children=False,
        sync_mode=ConnectorSyncMode.manual.value,
        sync_interval_minutes=None,
        status=ConnectorResourceStatus.active.value,
        next_auto_sync_at=None,
        provider_metadata={"upload_filename": filename},
        last_sync_started_at=utcnow(),
    )
    session.add(resource)
    await session.flush()

    sync_items = _iter_notion_export_items(filename, content_bytes)
    counts = await _ingest_sync_items_for_resource(
        session,
        connection=connection,
        resource=resource,
        sync_items=sync_items,
    )
    resource.last_sync_completed_at = utcnow()
    resource.last_sync_summary = counts
    await session.commit()
    await session.refresh(resource)
    return _resource_summary(resource)


async def update_resource(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    connection_id: UUID,
    resource_id: UUID,
    payload: ConnectorResourceUpdateRequest,
) -> ConnectorResourceSummary:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    if connection.owner_scope == ConnectorOwnerScope.workspace.value:
        _ensure_scope_permission(connection.owner_scope, auth_user)
    resource = await _get_resource_or_raise(session, resource_id)
    if resource.connection_id != connection.id:
        raise ConnectorNotFoundError("Connector resource does not belong to the connection.")
    if payload.visibility_scope is not None:
        resource.visibility_scope = _normalize_visibility_scope(connection.provider, resource.resource_kind, payload.visibility_scope)
    if payload.selection_mode is not None:
        resource.selection_mode = _normalize_selection_mode(connection.provider, resource.resource_kind, payload.selection_mode)
    resource.sync_children = _resource_sync_children_for_update(resource, payload.sync_children)
    if payload.status is not None:
        resource.status = payload.status
    if not _resource_supports_connector_sync(resource):
        if payload.sync_mode is not None and payload.sync_mode != ConnectorSyncMode.manual.value:
            raise ConnectorError(_non_syncable_resource_message(resource))
        if payload.sync_interval_minutes is not None:
            raise ConnectorError(_non_syncable_resource_message(resource))
        resource.sync_mode = ConnectorSyncMode.manual.value
        resource.sync_interval_minutes = None
        resource.next_auto_sync_at = None
    elif payload.sync_mode is not None or payload.sync_interval_minutes is not None:
        sync_mode, interval = _normalize_sync_schedule(
            payload.sync_mode or resource.sync_mode,
            payload.sync_interval_minutes if payload.sync_interval_minutes is not None else resource.sync_interval_minutes,
        )
        resource.sync_mode = sync_mode
        resource.sync_interval_minutes = interval
        resource.next_auto_sync_at = future_utc(seconds=interval * 60) if sync_mode == ConnectorSyncMode.auto.value and interval else None
    await session.commit()
    await session.refresh(resource)
    return _resource_summary(resource)


async def delete_resource(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    connection_id: UUID,
    resource_id: UUID,
) -> None:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    if connection.owner_scope == ConnectorOwnerScope.workspace.value:
        _ensure_scope_permission(connection.owner_scope, auth_user)
    resource = await _get_resource_or_raise(session, resource_id)
    if resource.connection_id != connection.id:
        raise ConnectorNotFoundError("Connector resource does not belong to the connection.")
    await session.delete(resource)
    await session.commit()


async def list_source_items(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    connection_id: UUID,
) -> list[ConnectorSourceItemSummary]:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    items = list(
        (
            await session.execute(
                select(ConnectorSourceItem)
                .where(ConnectorSourceItem.connection_id == connection.id)
                .order_by(ConnectorSourceItem.updated_at.desc())
                .limit(200)
            )
        ).scalars().all()
    )
    return [_source_item_summary(item) for item in items]


async def enqueue_connector_sync_job(
    session: AsyncSession,
    connection_id: UUID,
    resource_id: UUID,
    *,
    sync_mode: str,
    priority: int,
) -> ConnectorSyncJob:
    existing = (
        await session.execute(
            select(ConnectorSyncJob).where(
                ConnectorSyncJob.resource_id == resource_id,
                ConnectorSyncJob.status.in_([JobStatus.queued.value, JobStatus.processing.value]),
            )
        )
    ).scalar_one_or_none()
    if existing is not None:
        return existing
    job = ConnectorSyncJob(
        connection_id=connection_id,
        resource_id=resource_id,
        sync_mode=sync_mode,
        status=JobStatus.queued.value,
        priority=priority,
        payload={},
    )
    session.add(job)
    await session.flush()
    return job


async def request_resource_sync(
    session: AsyncSession,
    auth_user: AuthenticatedUser,
    connection_id: UUID,
    resource_id: UUID,
) -> JobSummary:
    connection = await _get_connection_or_raise(session, connection_id, auth_user)
    if connection.owner_scope == ConnectorOwnerScope.workspace.value:
        _ensure_scope_permission(connection.owner_scope, auth_user)
    resource = await _get_resource_or_raise(session, resource_id)
    if resource.connection_id != connection.id:
        raise ConnectorNotFoundError("Connector resource does not belong to the connection.")
    if not _resource_supports_connector_sync(resource):
        raise ConnectorError(_non_syncable_resource_message(resource))
    job = await enqueue_connector_sync_job(
        session,
        connection.id,
        resource.id,
        sync_mode=ConnectorSyncMode.manual.value,
        priority=80,
    )
    await session.commit()
    await session.refresh(job)
    return JobSummary(
        id=job.id,
        kind=ConnectorSyncJobKind.sync.value,
        title=f"리소스 동기화: {resource.name}",
        status=job.status,
        connection_id=job.connection_id,
        resource_id=job.resource_id,
        priority=job.priority,
        attempt_count=job.attempt_count,
        error_message=job.error_message,
        requested_at=job.requested_at,
        started_at=job.started_at,
        finished_at=job.finished_at,
    )


async def enqueue_due_sync_jobs(session: AsyncSession, *, limit: int = 10) -> int:
    due_resources = list(
        (
            await session.execute(
                select(ConnectorResource)
                .where(
                    ConnectorResource.sync_mode == ConnectorSyncMode.auto.value,
                    ConnectorResource.status == ConnectorResourceStatus.active.value,
                    ConnectorResource.next_auto_sync_at.is_not(None),
                    ConnectorResource.next_auto_sync_at <= utcnow(),
                )
                .order_by(ConnectorResource.next_auto_sync_at.asc())
                .limit(limit)
            )
        ).scalars().all()
    )
    created = 0
    corrected = 0
    for resource in due_resources:
        if not _resource_supports_connector_sync(resource):
            resource.sync_mode = ConnectorSyncMode.manual.value
            resource.sync_interval_minutes = None
            resource.next_auto_sync_at = None
            corrected += 1
            continue
        existing = (
            await session.execute(
                select(ConnectorSyncJob).where(
                    ConnectorSyncJob.resource_id == resource.id,
                    ConnectorSyncJob.status.in_([JobStatus.queued.value, JobStatus.processing.value]),
                )
            )
        ).scalar_one_or_none()
        if existing is not None:
            continue
        await enqueue_connector_sync_job(
            session,
            resource.connection_id,
            resource.id,
            sync_mode=ConnectorSyncMode.auto.value,
            priority=95,
        )
        created += 1
    if created or corrected:
        await session.commit()
    return created


async def acquire_next_connector_sync_job(session: AsyncSession) -> ConnectorSyncJob | None:
    job = (
        await session.execute(
            select(ConnectorSyncJob)
            .where(ConnectorSyncJob.status.in_([JobStatus.queued.value, JobStatus.failed.value]))
            .where(ConnectorSyncJob.attempt_count < get_settings().worker_max_attempts)
            .order_by(ConnectorSyncJob.priority.asc(), ConnectorSyncJob.requested_at.asc())
            .limit(1)
        )
    ).scalar_one_or_none()
    if job is None:
        return None
    job.status = JobStatus.processing.value
    job.started_at = utcnow()
    job.finished_at = None
    job.last_heartbeat_at = utcnow()
    job.attempt_count += 1
    job.error_message = None
    await session.flush()
    return job


async def mark_connector_job_failed(session: AsyncSession, job_id: UUID, message: str) -> None:
    job = await session.get(ConnectorSyncJob, job_id)
    if job is None:
        return
    job.status = JobStatus.failed.value
    job.error_message = message[:4000]
    job.finished_at = utcnow()
    job.last_heartbeat_at = utcnow()
    await session.flush()


async def mark_connector_job_completed(session: AsyncSession, job_id: UUID, payload: dict[str, Any]) -> None:
    job = await session.get(ConnectorSyncJob, job_id)
    if job is None:
        return
    job.status = JobStatus.completed.value
    job.error_message = None
    job.finished_at = utcnow()
    job.last_heartbeat_at = utcnow()
    job.payload = payload
    await session.flush()


async def _upsert_source_item(
    session: AsyncSession,
    *,
    connection_id: UUID,
    resource_id: UUID,
    item: PreparedSyncItem,
    document_id: UUID | None,
    status: str,
    unsupported_reason: str | None = None,
    error_message: str | None = None,
) -> ConnectorSourceItem:
    row = (
        await session.execute(
            select(ConnectorSourceItem).where(
                ConnectorSourceItem.connection_id == connection_id,
                ConnectorSourceItem.resource_id == resource_id,
                ConnectorSourceItem.external_item_id == item.external_item_id,
            )
        )
    ).scalar_one_or_none()
    if row is None:
        row = ConnectorSourceItem(
            connection_id=connection_id,
            resource_id=resource_id,
            external_item_id=item.external_item_id,
            mime_type=item.mime_type,
            name=item.title,
            source_url=item.source_url,
            source_revision_id=item.source_revision_id,
            internal_document_id=document_id,
            item_status=status,
            unsupported_reason=unsupported_reason,
            error_message=error_message,
            last_seen_at=utcnow(),
            last_synced_at=utcnow(),
            provider_metadata=item.provider_metadata,
        )
        session.add(row)
    else:
        row.mime_type = item.mime_type
        row.name = item.title
        row.source_url = item.source_url
        row.source_revision_id = item.source_revision_id
        row.internal_document_id = document_id
        row.item_status = status
        row.unsupported_reason = unsupported_reason
        row.error_message = error_message
        row.last_seen_at = utcnow()
        row.last_synced_at = utcnow()
        row.provider_metadata = item.provider_metadata
    await session.flush()
    return row


async def _archive_document_if_unreferenced(session: AsyncSession, document_id: UUID | None) -> None:
    if document_id is None:
        return
    remaining = int(
        (
            await session.execute(
                select(func.count(ConnectorSourceItem.id)).where(
                    ConnectorSourceItem.internal_document_id == document_id,
                    ConnectorSourceItem.item_status.in_(
                        [
                            ConnectorSourceItemStatus.imported.value,
                            ConnectorSourceItemStatus.unchanged.value,
                            ConnectorSourceItemStatus.failed.value,
                            ConnectorSourceItemStatus.unsupported.value,
                        ]
                    ),
                )
            )
        ).scalar_one()
    )
    if remaining > 0:
        return
    document = await session.get(Document, document_id)
    if document is not None:
        document.status = "archived"
        await session.flush()


async def _ingest_sync_items_for_resource(
    session: AsyncSession,
    *,
    connection: ConnectorConnection,
    resource: ConnectorResource,
    sync_items: list[PreparedSyncItem],
) -> dict[str, int]:
    seen_ids = {item.external_item_id for item in sync_items}
    counts = {"imported": 0, "unchanged": 0, "unsupported": 0, "failed": 0, "deleted": 0}

    for item in sync_items:
        try:
            if item.unsupported_reason:
                counts["unsupported"] += 1
                await _upsert_source_item(
                    session,
                    connection_id=connection.id,
                    resource_id=resource.id,
                    item=item,
                    document_id=None,
                    status=ConnectorSourceItemStatus.unsupported.value,
                    unsupported_reason=item.unsupported_reason,
                )
                await session.commit()
                continue

            if not item.content or not item.content.strip():
                raise ConnectorError("Extracted content is empty.")

            payload = IngestDocumentRequest(
                source_system=(
                    "google-drive"
                    if connection.provider == ConnectorProvider.google_drive.value
                    else "github"
                    if connection.provider == ConnectorProvider.github.value
                    else "notion-export"
                    if resource.selection_mode == "export_upload"
                    else "notion"
                ),
                source_external_id=item.external_item_id,
                source_revision_id=item.source_revision_id,
                source_url=item.source_url,
                slug=_stable_document_slug(item.title, item.external_item_id),
                title=item.title,
                content_type=item.content_type,  # type: ignore[arg-type]
                content=item.content,
                doc_type=item.doc_type,
                language_code="ko",
                status="published",
                visibility_scope=resource.visibility_scope,  # type: ignore[arg-type]
                metadata={
                    "provider": connection.provider,
                    "connector_connection_id": str(connection.id),
                    "connector_resource_id": str(resource.id),
                    "selection_mode": resource.selection_mode,
                    "visibility_scope": resource.visibility_scope,
                    **item.provider_metadata,
                },
                priority=110,
            )
            result = await ingest_document(
                session,
                payload,
                workspace_id=connection.workspace_id,
            )
            status = ConnectorSourceItemStatus.unchanged.value if result.unchanged else ConnectorSourceItemStatus.imported.value
            counts["unchanged" if result.unchanged else "imported"] += 1
            await _upsert_source_item(
                session,
                connection_id=connection.id,
                resource_id=resource.id,
                item=item,
                document_id=result.document.id,
                status=status,
            )
        except ConnectorError as exc:
            counts["unsupported"] += 1
            await _upsert_source_item(
                session,
                connection_id=connection.id,
                resource_id=resource.id,
                item=item,
                document_id=None,
                status=ConnectorSourceItemStatus.unsupported.value,
                unsupported_reason=str(exc),
            )
            await session.commit()
        except Exception as exc:  # noqa: BLE001
            counts["failed"] += 1
            await _upsert_source_item(
                session,
                connection_id=connection.id,
                resource_id=resource.id,
                item=item,
                document_id=None,
                status=ConnectorSourceItemStatus.failed.value,
                error_message=str(exc),
            )
            await session.commit()

    existing_items = list(
        (
            await session.execute(
                select(ConnectorSourceItem).where(ConnectorSourceItem.resource_id == resource.id)
            )
        ).scalars().all()
    )
    for item in existing_items:
        if item.external_item_id in seen_ids:
            continue
        item.item_status = ConnectorSourceItemStatus.deleted.value
        item.last_synced_at = utcnow()
        counts["deleted"] += 1
        await _archive_document_if_unreferenced(session, item.internal_document_id)

    return counts


async def process_connector_sync_job(session_factory: async_sessionmaker[AsyncSession], job_id: UUID) -> None:
    async with session_factory() as session:
        job = await session.get(ConnectorSyncJob, job_id)
        if job is None:
            return
        connection = await session.get(ConnectorConnection, job.connection_id)
        resource = await session.get(ConnectorResource, job.resource_id)
        if connection is None or resource is None:
            raise ConnectorError("Connector sync job references missing connection/resource.")
        resource.last_sync_started_at = utcnow()
        await session.commit()

        sync_items = await _sync_items_for_resource(session, connection, resource)
        counts = await _ingest_sync_items_for_resource(
            session,
            connection=connection,
            resource=resource,
            sync_items=sync_items,
        )

        resource.last_sync_completed_at = utcnow()
        resource.last_sync_summary = counts
        if resource.sync_mode == ConnectorSyncMode.auto.value and resource.sync_interval_minutes:
            resource.next_auto_sync_at = future_utc(seconds=resource.sync_interval_minutes * 60)
        else:
            resource.next_auto_sync_at = None
        await mark_connector_job_completed(session, job.id, counts)
        await session.commit()
